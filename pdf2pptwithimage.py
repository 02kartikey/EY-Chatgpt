from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader
import openai
import fitz  # PyMuPDF library

import os


# Function to extract images from a specific page of the PDF
def extract_images_from_pdf(pdf_file, page_number):
    images = []
    pdf_document = fitz.open(pdf_file)
    page = pdf_document[page_number - 1]
    image_list = page.get_images(full=True)

    # Get the full path of the "images" directory based on the current working directory
    current_directory = os.getcwd()
    images_directory = os.path.join(current_directory, "images")

    # Create the "images" directory if it doesn't exist
    if not os.path.exists(images_directory):
        os.makedirs(images_directory)

    for img_index, img in enumerate(image_list):
        image = pdf_document.extract_image(img[0])
        image_name = f"image_page{page_number}_{img_index + 1}.png"
        image_path = os.path.join(images_directory, image_name)  # Save images in the "images" subdirectory
        with open(image_path, "wb") as img_file:
            img_file.write(image["image"])
        images.append(image_path)

    return images

# Rest of the code remains the same...


# Function to extract text from a specific page of the PDF
def extract_text_from_pdf(pdf_file, page_number):
    with open(pdf_file, 'rb') as file:
        reader = PdfReader(file)
        page = reader.pages[page_number - 1]
        text = page.extract_text()
    return text

# Function to summarize the text using OpenAI API
def summarize_with_openai(text):
    # Replace with your actual OpenAI API key
    openai.api_key = "OpenAI-Key"
    response = openai.Completion.create(
        engine="text-davinci-002",
        prompt=text,
        temperature=0.7,
        max_tokens=150
    )
    summary = response['choices'][0]['text']
    return summary

# Function to extract images from a specific page of the PDF
def extract_images_from_pdf(pdf_file, page_number):
    images = []
    pdf_document = fitz.open(pdf_file)
    page = pdf_document[page_number - 1]
    image_list = page.get_images(full=True)

    for img_index, img in enumerate(image_list):
        image = pdf_document.extract_image(img[0])
        image_name = f"image_page{page_number}_{img_index + 1}.png"
        image_path = f"images/{image_name}"  # Save images in a subdirectory named "images"
        with open(image_path, "wb") as img_file:
            img_file.write(image["image"])
        images.append(image_path)

    return images

## ... (previous code)

def main():
    pdf_file = '####.pdf'  # Replace with the path to your PDF file
    total_pages = 0
    summaries = []
    all_images = []

    with open(pdf_file, 'rb') as file:
        reader = PdfReader(file)
        total_pages = len(reader.pages)

    print(f"Total Pages in the PDF: {total_pages}\n")

    # Extract all images once before the loop
    for page_number in range(1, total_pages + 1):
        images = extract_images_from_pdf(pdf_file, page_number)
        all_images.append(images)

    for page_number in range(1, total_pages + 1):
        # Step 1: Extract the text from the PDF page
        extracted_text = extract_text_from_pdf(pdf_file, page_number)

        # Step 2: Summarize the extracted text using OpenAI API
        summary = summarize_with_openai(extracted_text)

        # Append the summary to the list
        summaries.append(summary)

        # Step 3: Extracted images from the PDF page (already done before the loop)
        images = all_images[page_number - 1]

        # Display the summarized output for each page (optional)
        print(f"Summarized Output for Page {page_number}:")
        print(summary)
        print("----------------------\n")

    # Create a new presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Choose a slide layout (Title and Content)

    # Add a slide for each summary
    for page_number, (summary, images) in enumerate(zip(summaries, all_images), start=1):
        slide = prs.slides.add_slide(slide_layout)

        # Set the title text
        title = slide.shapes.title
        title.text = f"Here is your {page_number} Summary"

        # Set the content text only if the layout has a content placeholder (index 1)
        content = slide.placeholders[1]
        if content.placeholder_format.idx == 1:
            content.text = summary

        # Add the images to the slide
        for img_path in images:
            slide.shapes.add_picture(img_path, Inches(2), Inches(2), height=Inches(4.5))

    # Save the presentation to a PPT file
    prs.save("summary_presentation_with_image.pptx")

    print("Presentation with text summaries and images saved successfully!")

# ... (rest of the code remains the same)

if __name__ == "__main__":
    main()
