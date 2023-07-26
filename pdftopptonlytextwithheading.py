from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader
import openai

# Function to extract text from a specific page of the PDF
def extract_text_from_pdf(pdf_file, page_number):
    with open(pdf_file, 'rb') as file:
        reader = PdfReader(file)
        page = reader.pages[page_number - 1]
        text = page.extract_text()
    return text

# Function to summarize the text using OpenAI API
def summarize_with_openai(text):
    # Replace with your actual OpenAI API key
    openai.api_key = "OpenAPI KEY"
    response = openai.Completion.create(
        engine="text-davinci-002",
        prompt=text,
        temperature=0.7,
        max_tokens=150
    )
    summary = response['choices'][0]['text']
    return summary

def main():
    pdf_file = '#######.pdf'  # Replace with the path to your PDF file
    total_pages = 0
    summaries = []

    with open(pdf_file, 'rb') as file:
        reader = PdfReader(file)
        total_pages = len(reader.pages)

        # Extract the title from the first page of the PDF
        title_text = extract_text_from_pdf(pdf_file, 1)

    print(f"Total Pages in the PDF: {total_pages}\n")

    for page_number in range(1, total_pages + 1):
        # Step 1: Extract the text from the PDF page
        extracted_text = extract_text_from_pdf(pdf_file, page_number)

        # Step 2: Summarize the extracted text using OpenAI API
        summary = summarize_with_openai(extracted_text)

        # Append the summary to the list
        summaries.append(summary)

        # Display the summarized output for each page (optional)
        print(f"Summarized Output for Page {page_number}:")
        print(summary)
        print("----------------------\n")

    # Create a new presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Choose a slide layout (Title Slide with Content)

    # Add a slide for each summary
    for page_number, summary in enumerate(summaries, start=1):
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"{title_text} - Page {page_number} Summary"
        content = slide.placeholders[1]
        content.text = summary

    # Save the presentation to a PPT file
    prs.save("summary_presentation.pptx")

    print("Presentation with text summaries saved successfully!")

if __name__ == "__main__":
    main()
